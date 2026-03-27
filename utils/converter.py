from pdf2docx import Converter
import subprocess
import os
import mammoth
import comtypes.client
import pythoncom


# ===============================
# DOCX → HTML
# ===============================
def convert_docx_to_html(input_path):
    with open(input_path, "rb") as docx_file:
        result = mammoth.convert_to_html(docx_file)
        return result.value


# ===============================
# PDF → DOCX
# ===============================
def convert_pdf_to_docx(input_path, output_path):
    cv = Converter(input_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()


# ===============================
# DOCX → PDF (LibreOffice)
# ===============================
def convert_docx_to_pdf(input_path, output_folder):

    input_path = os.path.abspath(input_path)
    output_folder = os.path.abspath(output_folder)

    command = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        "--headless",
        "--convert-to",
        "pdf",
        input_path,
        "--outdir",
        output_folder
    ]

    result = subprocess.run(command, capture_output=True, text=True)

    if result.returncode != 0:
        print(result.stderr)
        raise Exception("DOCX to PDF conversion failed.")


# ===============================
# WORD → PDF (MS Word COM)
# ===============================
def convert_word_to_pdf(input_path, output_path):

    pythoncom.CoInitialize()

    try:
        word = comtypes.client.CreateObject("Word.Application")
        word.Visible = False

        doc = word.Documents.Open(os.path.abspath(input_path))
        doc.SaveAs(os.path.abspath(output_path), FileFormat=17)
        doc.Close()
        word.Quit()

    finally:
        pythoncom.CoUninitialize()


# ===============================
# POWERPOINT → PDF
# ===============================
def convert_ppt_to_pdf(input_path, output_path):

    pythoncom.CoInitialize()

    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1

        presentation = powerpoint.Presentations.Open(os.path.abspath(input_path))
        presentation.SaveAs(os.path.abspath(output_path), 32)
        presentation.Close()

        powerpoint.Quit()

    finally:
        pythoncom.CoUninitialize()


# ===============================
# EXCEL → PDF
# ===============================
def convert_excel_to_pdf(input_path, output_path):

    pythoncom.CoInitialize()

    try:
        excel = comtypes.client.CreateObject("Excel.Application")
        excel.Visible = False

        workbook = excel.Workbooks.Open(os.path.abspath(input_path))
        workbook.ExportAsFixedFormat(0, os.path.abspath(output_path))
        workbook.Close()

        excel.Quit()

    finally:
        pythoncom.CoUninitialize()

from pptx import Presentation

def create_ppt_from_ai(slides, output_path):

    prs = Presentation()

    for slide_data in slides:

        title = slide_data.get("title", "")
        points = slide_data.get("points", [])

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title

        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        for point in points:
            p = body.add_paragraph()
            p.text = point
            p.level = 0

    prs.save(output_path)        


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def create_ppt_from_ai(slides, output_path, theme="modern"):

    prs = Presentation()

    for slide_data in slides:

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content = slide.placeholders[1]

        # 🎯 TITLE
        title.text = slide_data.get("title", "")

        # 🎯 CONTENT
        points = slide_data.get("points", [])
        tf = content.text_frame
        tf.clear()

        for i, point in enumerate(points):
            if i == 0:
                tf.text = point
            else:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

        # 🎯 SPEAKER NOTES
        notes = slide_data.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        # 🎨 THEME STYLING
        if theme == "dark":
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(20, 20, 20)

            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        elif theme == "business":
            title.text_frame.paragraphs[0].font.size = Pt(32)

        elif theme == "modern":
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True

    prs.save(output_path)